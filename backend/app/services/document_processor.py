"""Document processing and text extraction with OCR support"""
import os
import io
from pathlib import Path
from typing import Dict, Any, Optional, List
import asyncio
from datetime import datetime

# Document processing libraries
import markdown
from bs4 import BeautifulSoup
import html2text


class DocumentProcessor:
    """Process various document formats and extract text with OCR capabilities"""

    SUPPORTED_EXTENSIONS = {
        # Text formats
        '.txt': 'text',
        '.md': 'markdown',
        '.markdown': 'markdown',
        '.html': 'html',
        '.htm': 'html',
        '.rtf': 'rtf',
        # Documents
        '.pdf': 'pdf',
        '.docx': 'docx',
        '.doc': 'docx',
        '.xlsx': 'excel',
        '.xls': 'excel',
        '.pptx': 'powerpoint',
        '.ppt': 'powerpoint',
        # E-books
        '.epub': 'epub',
        # Images (OCR)
        '.png': 'image',
        '.jpg': 'image',
        '.jpeg': 'image',
        '.gif': 'image',
        '.bmp': 'image',
        '.tiff': 'image',
        '.tif': 'image',
        '.webp': 'image',
        # Code files
        '.py': 'code',
        '.js': 'code',
        '.ts': 'code',
        '.jsx': 'code',
        '.tsx': 'code',
        '.java': 'code',
        '.cpp': 'code',
        '.c': 'code',
        '.h': 'code',
        '.css': 'code',
        '.scss': 'code',
        '.json': 'code',
        '.yaml': 'code',
        '.yml': 'code',
        '.xml': 'code',
        '.sql': 'code',
        '.sh': 'code',
        '.bash': 'code',
        '.rs': 'code',
        '.go': 'code',
        '.rb': 'code',
        '.php': 'code',
        '.swift': 'code',
        '.kt': 'code',
        '.scala': 'code',
        '.r': 'code',
        '.vue': 'code',
        '.svelte': 'code',
        # Data files
        '.csv': 'data',
        '.tsv': 'data',
    }

    # OCR available flag
    _ocr_available: Optional[bool] = None

    def __init__(self):
        self.html_converter = html2text.HTML2Text()
        self.html_converter.ignore_links = False
        self.html_converter.ignore_images = True
        self.html_converter.body_width = 0
        self._check_ocr_availability()

    def _check_ocr_availability(self):
        """Check if OCR libraries are available"""
        if self._ocr_available is None:
            try:
                import pytesseract
                from PIL import Image
                self._ocr_available = True
            except ImportError:
                self._ocr_available = False
                print("OCR not available. Install pytesseract and Pillow for image text extraction.")

    @property
    def ocr_available(self) -> bool:
        """Check if OCR is available"""
        if self._ocr_available is None:
            self._check_ocr_availability()
        return self._ocr_available

    async def process_file(self, path: Path) -> Dict[str, Any]:
        """
        Process a file and extract its content.

        Returns dict with:
            - content: extracted text content
            - file_type: type category
            - title: document title if found
            - metadata: additional metadata
        """
        ext = path.suffix.lower()
        file_type = self.SUPPORTED_EXTENSIONS.get(ext, 'unknown')

        if file_type == 'unknown':
            # Try to read as text anyway
            try:
                content = await self._read_text_file(path)
                return {
                    'content': content,
                    'file_type': 'text',
                    'title': path.stem,
                    'metadata': {'original_extension': ext}
                }
            except Exception:
                raise ValueError(f"Unsupported file type: {ext}")

        # Route to appropriate processor
        processors = {
            'text': self._process_text,
            'markdown': self._process_markdown,
            'html': self._process_html,
            'pdf': self._process_pdf,
            'docx': self._process_docx,
            'code': self._process_code,
            'data': self._process_data,
            'excel': self._process_excel,
            'powerpoint': self._process_powerpoint,
            'epub': self._process_epub,
            'image': self._process_image,
            'rtf': self._process_rtf,
        }

        processor = processors.get(file_type, self._process_text)
        return await processor(path)

    async def _read_text_file(self, path: Path) -> str:
        """Read a text file with encoding detection"""
        encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']

        for encoding in encodings:
            try:
                return path.read_text(encoding=encoding)
            except (UnicodeDecodeError, LookupError):
                continue

        raise ValueError(f"Could not decode file: {path}")

    async def _process_text(self, path: Path) -> Dict[str, Any]:
        """Process plain text file"""
        content = await self._read_text_file(path)
        return {
            'content': content,
            'file_type': 'text',
            'title': path.stem,
            'metadata': {}
        }

    async def _process_markdown(self, path: Path) -> Dict[str, Any]:
        """Process Markdown file"""
        content = await self._read_text_file(path)

        # Extract title from first heading if present
        title = path.stem
        lines = content.split('\n')
        for line in lines:
            if line.startswith('# '):
                title = line[2:].strip()
                break

        # Convert to HTML then to plain text for better chunking
        html = markdown.markdown(content)
        plain_text = self.html_converter.handle(html)

        return {
            'content': plain_text,
            'file_type': 'markdown',
            'title': title,
            'metadata': {'has_frontmatter': content.startswith('---')}
        }

    async def _process_html(self, path: Path) -> Dict[str, Any]:
        """Process HTML file"""
        content = await self._read_text_file(path)
        soup = BeautifulSoup(content, 'html.parser')

        # Extract title
        title = path.stem
        title_tag = soup.find('title')
        if title_tag:
            title = title_tag.get_text().strip()

        # Remove script and style elements
        for element in soup(['script', 'style', 'nav', 'footer', 'header']):
            element.decompose()

        # Get main content
        main_content = soup.find('main') or soup.find('article') or soup.find('body') or soup

        plain_text = self.html_converter.handle(str(main_content))

        return {
            'content': plain_text,
            'file_type': 'html',
            'title': title,
            'metadata': {}
        }

    async def _process_pdf(self, path: Path) -> Dict[str, Any]:
        """Process PDF file"""
        try:
            from PyPDF2 import PdfReader

            reader = PdfReader(str(path))
            text_parts = []

            for page in reader.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)

            content = '\n\n'.join(text_parts)

            # Try to get title from metadata
            title = path.stem
            if reader.metadata and reader.metadata.title:
                title = reader.metadata.title

            return {
                'content': content,
                'file_type': 'pdf',
                'title': title,
                'metadata': {
                    'page_count': len(reader.pages),
                    'author': reader.metadata.author if reader.metadata else None
                }
            }
        except Exception as e:
            raise ValueError(f"Could not process PDF: {e}")

    async def _process_docx(self, path: Path) -> Dict[str, Any]:
        """Process Word document"""
        try:
            from docx import Document

            doc = Document(str(path))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            content = '\n\n'.join(paragraphs)

            # Try to get title from properties
            title = path.stem
            if doc.core_properties.title:
                title = doc.core_properties.title

            return {
                'content': content,
                'file_type': 'docx',
                'title': title,
                'metadata': {
                    'author': doc.core_properties.author,
                    'created': doc.core_properties.created.isoformat() if doc.core_properties.created else None
                }
            }
        except Exception as e:
            raise ValueError(f"Could not process DOCX: {e}")

    async def _process_code(self, path: Path) -> Dict[str, Any]:
        """Process code file"""
        content = await self._read_text_file(path)

        # Add file info as header
        header = f"File: {path.name}\nLanguage: {path.suffix[1:]}\n\n"

        return {
            'content': header + content,
            'file_type': 'code',
            'title': path.name,
            'metadata': {
                'language': path.suffix[1:],
                'line_count': len(content.split('\n'))
            }
        }

    async def _process_data(self, path: Path) -> Dict[str, Any]:
        """Process data files (CSV, TSV, etc.)"""
        content = await self._read_text_file(path)

        # For data files, we keep the raw content but note it's structured
        return {
            'content': content,
            'file_type': 'data',
            'title': path.name,
            'metadata': {
                'row_count': len(content.split('\n')),
                'format': path.suffix[1:]
            }
        }

    def process_web_content(self, html: str, url: str, title: Optional[str] = None) -> Dict[str, Any]:
        """Process captured web content"""
        soup = BeautifulSoup(html, 'html.parser')

        # Extract title
        if not title:
            title_tag = soup.find('title')
            title = title_tag.get_text().strip() if title_tag else url

        # Remove unwanted elements
        for element in soup(['script', 'style', 'nav', 'footer', 'header',
                           'aside', 'iframe', 'noscript', 'form']):
            element.decompose()

        # Try to find main content area
        main_content = (
            soup.find('article') or
            soup.find('main') or
            soup.find(class_=['content', 'article', 'post', 'entry']) or
            soup.find('body') or
            soup
        )

        plain_text = self.html_converter.handle(str(main_content))

        # Clean up excessive whitespace
        lines = [line.strip() for line in plain_text.split('\n')]
        lines = [line for line in lines if line]
        clean_text = '\n\n'.join(lines)

        return {
            'content': clean_text,
            'title': title,
            'url': url
        }

    async def _process_excel(self, path: Path) -> Dict[str, Any]:
        """Process Excel files (.xlsx, .xls)"""
        try:
            import openpyxl

            wb = openpyxl.load_workbook(str(path), data_only=True)
            content_parts = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                content_parts.append(f"## Sheet: {sheet_name}\n")

                rows = []
                for row in sheet.iter_rows(values_only=True):
                    row_text = '\t'.join([str(cell) if cell is not None else '' for cell in row])
                    if row_text.strip():
                        rows.append(row_text)

                content_parts.append('\n'.join(rows))
                content_parts.append('\n')

            content = '\n'.join(content_parts)

            return {
                'content': content,
                'file_type': 'excel',
                'title': path.stem,
                'metadata': {
                    'sheet_count': len(wb.sheetnames),
                    'sheet_names': wb.sheetnames
                }
            }
        except Exception as e:
            raise ValueError(f"Could not process Excel file: {e}")

    async def _process_powerpoint(self, path: Path) -> Dict[str, Any]:
        """Process PowerPoint files (.pptx)"""
        try:
            from pptx import Presentation

            prs = Presentation(str(path))
            content_parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                content_parts.append(f"\n## Slide {slide_num}\n")

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content_parts.append(shape.text)

            content = '\n'.join(content_parts)

            return {
                'content': content,
                'file_type': 'powerpoint',
                'title': path.stem,
                'metadata': {
                    'slide_count': len(prs.slides)
                }
            }
        except Exception as e:
            raise ValueError(f"Could not process PowerPoint file: {e}")

    async def _process_epub(self, path: Path) -> Dict[str, Any]:
        """Process EPUB e-book files"""
        try:
            import ebooklib
            from ebooklib import epub

            book = epub.read_epub(str(path))
            content_parts = []

            # Get title
            title = path.stem
            if book.get_metadata('DC', 'title'):
                title = book.get_metadata('DC', 'title')[0][0]

            # Extract text from all document items
            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    soup = BeautifulSoup(item.get_content(), 'html.parser')
                    text = soup.get_text(separator='\n')
                    if text.strip():
                        content_parts.append(text)

            content = '\n\n'.join(content_parts)

            # Get author
            author = None
            if book.get_metadata('DC', 'creator'):
                author = book.get_metadata('DC', 'creator')[0][0]

            return {
                'content': content,
                'file_type': 'epub',
                'title': title,
                'metadata': {
                    'author': author,
                    'chapter_count': len(content_parts)
                }
            }
        except Exception as e:
            raise ValueError(f"Could not process EPUB file: {e}")

    async def _process_image(self, path: Path) -> Dict[str, Any]:
        """Process image files using OCR"""
        if not self.ocr_available:
            return {
                'content': f"[Image file: {path.name}. OCR not available - install pytesseract]",
                'file_type': 'image',
                'title': path.stem,
                'metadata': {'ocr_performed': False}
            }

        try:
            import pytesseract
            from PIL import Image

            image = Image.open(str(path))

            # Perform OCR
            text = pytesseract.image_to_string(image)

            return {
                'content': text.strip() or "[No text detected in image]",
                'file_type': 'image',
                'title': path.stem,
                'metadata': {
                    'ocr_performed': True,
                    'image_size': f"{image.width}x{image.height}",
                    'image_mode': image.mode
                }
            }
        except Exception as e:
            raise ValueError(f"Could not process image file: {e}")

    async def _process_rtf(self, path: Path) -> Dict[str, Any]:
        """Process RTF files"""
        try:
            from striprtf.striprtf import rtf_to_text

            rtf_content = path.read_bytes().decode('utf-8', errors='ignore')
            content = rtf_to_text(rtf_content)

            return {
                'content': content,
                'file_type': 'rtf',
                'title': path.stem,
                'metadata': {}
            }
        except Exception as e:
            raise ValueError(f"Could not process RTF file: {e}")

    async def process_pdf_with_ocr(self, path: Path) -> Dict[str, Any]:
        """Process PDF with OCR fallback for scanned documents"""
        try:
            from PyPDF2 import PdfReader

            reader = PdfReader(str(path))
            text_parts = []
            used_ocr = False

            for page_num, page in enumerate(reader.pages):
                text = page.extract_text()

                # If text extraction yields little/no text, try OCR
                if (not text or len(text.strip()) < 50) and self.ocr_available:
                    ocr_text = await self._ocr_pdf_page(path, page_num)
                    if ocr_text:
                        text = ocr_text
                        used_ocr = True

                if text:
                    text_parts.append(text)

            content = '\n\n'.join(text_parts)

            title = path.stem
            if reader.metadata and reader.metadata.title:
                title = reader.metadata.title

            return {
                'content': content,
                'file_type': 'pdf',
                'title': title,
                'metadata': {
                    'page_count': len(reader.pages),
                    'author': reader.metadata.author if reader.metadata else None,
                    'used_ocr': used_ocr
                }
            }
        except Exception as e:
            raise ValueError(f"Could not process PDF: {e}")

    async def _ocr_pdf_page(self, pdf_path: Path, page_num: int) -> Optional[str]:
        """OCR a single PDF page"""
        try:
            import pdf2image
            import pytesseract

            # Convert PDF page to image
            images = pdf2image.convert_from_path(
                str(pdf_path),
                first_page=page_num + 1,
                last_page=page_num + 1,
                dpi=300
            )

            if images:
                text = pytesseract.image_to_string(images[0])
                return text.strip()
        except Exception:
            pass
        return None

    async def summarize_document(self, content: str, title: str = "") -> Dict[str, Any]:
        """Generate a summary of document content"""
        from .ollama_service import ollama_service
        from ..core.config import settings

        # Truncate if too long
        max_chars = 15000
        truncated = len(content) > max_chars
        content_to_summarize = content[:max_chars] if truncated else content

        prompt = f"""Please provide a concise summary of the following document.
Include:
1. Main topic/purpose
2. Key points (3-5 bullets)
3. Any important conclusions or action items

Document title: {title}

Content:
{content_to_summarize}

{'[Document truncated due to length]' if truncated else ''}

Summary:"""

        try:
            result = await ollama_service.generate(
                model=settings.MODELS['document'],
                prompt=prompt,
                options={"temperature": 0.3}
            )
            summary = result.get("response", "")

            return {
                'summary': summary,
                'truncated': truncated,
                'original_length': len(content)
            }
        except Exception as e:
            return {
                'summary': f"Could not generate summary: {e}",
                'truncated': truncated,
                'original_length': len(content)
            }

    def get_supported_extensions(self) -> List[str]:
        """Get list of supported file extensions"""
        return list(self.SUPPORTED_EXTENSIONS.keys())

    def get_file_type(self, path: Path) -> str:
        """Get the file type category for a path"""
        ext = path.suffix.lower()
        return self.SUPPORTED_EXTENSIONS.get(ext, 'unknown')


# Singleton instance
document_processor = DocumentProcessor()
